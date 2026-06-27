# backend/exports/ppt_export.py
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import Dict, Any, List

def add_header(slide, title_text: str):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = "Segoe UI"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42) # Slate-900

def set_background(slide, color):
    # Sets background color via solid fill
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def generate_ppt(session_id: str, filename: str, dataset_type: str, kpis: List[Dict[str, Any]], insights: Dict[str, Any]) -> io.BytesIO:
    prs = Presentation()
    # Use widescreen layout 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Colors
    slate_navy = RGBColor(15, 23, 42) # #0F172A
    bright_blue = RGBColor(29, 78, 216) # #1D4ED8
    pure_white = RGBColor(255, 255, 255)
    soft_gray = RGBColor(241, 245, 249) # #F1F5F9
    
    # Blank slide layout is index 6
    blank_layout = prs.slide_layouts[6]
    
    # ----------------------------------------------------
    # Slide 1: Title Slide (Dark Background)
    # ----------------------------------------------------
    slide = prs.slides.add_slide(blank_layout)
    set_background(slide, slate_navy)
    
    # Title Text Frame
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(8.0), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "BOARDROOM-AI"
    p.font.name = "Segoe UI"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = bright_blue
    p.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = f"{dataset_type} Executive Briefing"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = pure_white
    p2.space_after = Pt(20)
    
    p3 = tf.add_paragraph()
    health_score = insights.get("health_score", 8.0)
    p3.text = f"Dataset: {filename}   |   Business Health Score: {health_score}/10"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(148, 163, 184) # Muted gray
    
    # Speaker Notes
    slide.notes_slide.notes_text_frame.text = (
        f"Welcome to the {dataset_type} Executive Briefing deck. "
        f"This report covers insights generated from {filename}. "
        f"Overall corporate health index scored at {health_score} out of 10. "
        "Today we will walk through the core performance drivers, concerns, and strategic recommendations."
    )
    
    # ----------------------------------------------------
    # Slide 2: Executive Summary (White Background)
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_background(slide2, pure_white)
    add_header(slide2, "Executive Summary")
    
    # Left column: Paragraph
    summary_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.0), Inches(3.8))
    tf_sum = summary_box.text_frame
    tf_sum.word_wrap = True
    
    p_sum_title = tf_sum.paragraphs[0]
    p_sum_title.text = "Operational Overview"
    p_sum_title.font.name = "Segoe UI"
    p_sum_title.font.size = Pt(18)
    p_sum_title.font.bold = True
    p_sum_title.font.color.rgb = bright_blue
    p_sum_title.space_after = Pt(10)
    
    p_sum_body = tf_sum.add_paragraph()
    p_sum_body.text = insights.get("executive_summary", "")
    p_sum_body.font.name = "Segoe UI"
    p_sum_body.font.size = Pt(12)
    p_sum_body.font.color.rgb = slate_navy
    p_sum_body.line_spacing = 1.3
    
    # Right column: Health Gauge Box
    gauge_box = slide2.shapes.add_textbox(Inches(6.0), Inches(1.3), Inches(3.5), Inches(3.8))
    tf_gauge = gauge_box.text_frame
    tf_gauge.word_wrap = True
    
    p_g_title = tf_gauge.paragraphs[0]
    p_g_title.text = "Business Health Rating"
    p_g_title.alignment = PP_ALIGN.CENTER
    p_g_title.font.name = "Segoe UI"
    p_g_title.font.size = Pt(18)
    p_g_title.font.bold = True
    p_g_title.font.color.rgb = slate_navy
    p_g_title.space_after = Pt(20)
    
    p_g_num = tf_gauge.add_paragraph()
    p_g_num.text = f"{health_score} / 10"
    p_g_num.alignment = PP_ALIGN.CENTER
    p_g_num.font.name = "Segoe UI"
    p_g_num.font.size = Pt(54)
    p_g_num.font.bold = True
    p_g_num.font.color.rgb = bright_blue
    p_g_num.space_after = Pt(10)
    
    # Speaker Notes
    slide2.notes_slide.notes_text_frame.text = (
        f"Here we highlight the core executive summary. "
        "To summarize the current operational status, the business scores a "
        f"{health_score}/10 based on standard statistical indicators. "
        "Let's move onto the details."
    )
    
    # ----------------------------------------------------
    # Slide 3: KPI Dashboard Grid (White Background)
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_background(slide3, pure_white)
    add_header(slide3, "Key Performance Indicators")
    
    # Draw a 2x3 grid of textboxes for KPIs
    rows = 2
    cols = 3
    width = Inches(2.8)
    height = Inches(1.5)
    
    for idx, kpi in enumerate(kpis[:6]):
        r = idx // cols
        c = idx % cols
        
        left = Inches(0.5 + c * 3.1)
        top = Inches(1.3 + r * 1.8)
        
        # Add KPI Card
        card = slide3.shapes.add_textbox(left, top, width, height)
        tf_card = card.text_frame
        tf_card.word_wrap = True
        
        # Style background or border via shape is hard for textboxes, so we just write styled text
        p_name = tf_card.paragraphs[0]
        p_name.text = kpi.get("name").upper()
        p_name.font.name = "Segoe UI"
        p_name.font.size = Pt(11)
        p_name.font.bold = True
        p_name.font.color.rgb = RGBColor(100, 116, 139) # slate-500
        p_name.space_after = Pt(5)
        
        p_val = tf_card.add_paragraph()
        p_val.text = kpi.get("value")
        p_val.font.name = "Segoe UI"
        p_val.font.size = Pt(28)
        p_val.font.bold = True
        p_val.font.color.rgb = bright_blue
        p_val.space_after = Pt(5)
        
        p_ins = tf_card.add_paragraph()
        p_ins.text = kpi.get("insight")
        p_ins.font.name = "Segoe UI"
        p_ins.font.size = Pt(9.5)
        p_ins.font.color.rgb = slate_navy
        
    # Speaker Notes
    slide3.notes_slide.notes_text_frame.text = (
        "This dashboard grid displays our core calculated metrics. "
        "Specifically, " + ", ".join([f"{k['name']} is {k['value']}" for k in kpis[:4]]) + "."
    )
    
    # ----------------------------------------------------
    # Slide 4: Detailed Insights List (White Background)
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_background(slide4, pure_white)
    add_header(slide4, "Analytics & Statistical Insights")
    
    ins_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(3.9))
    tf_ins = ins_box.text_frame
    tf_ins.word_wrap = True
    
    for idx, ins in enumerate(insights.get("insights", [])[:4]):
        p_ins = tf_ins.paragraphs[0] if idx == 0 else tf_ins.add_paragraph()
        p_ins.text = f"• {ins.get('category').upper()}: {ins.get('headline')}"
        p_ins.font.name = "Segoe UI"
        p_ins.font.size = Pt(14)
        p_ins.font.bold = True
        p_ins.font.color.rgb = bright_blue
        p_ins.space_before = Pt(8)
        p_ins.space_after = Pt(3)
        
        p_ins_desc = tf_ins.add_paragraph()
        p_ins_desc.text = f"  {ins.get('detail')} (Data Check: {ins.get('data_point')})"
        p_ins_desc.font.name = "Segoe UI"
        p_ins_desc.font.size = Pt(11)
        p_ins_desc.font.color.rgb = slate_navy
        
    # Speaker Notes
    slide4.notes_slide.notes_text_frame.text = (
        "Reviewing the statistical insights: "
        "We see a clear distribution and segmentation breakdown. "
        "The correlation matrix and clustering algorithms highlight distinct groupings in our user populations."
    )
    
    # ----------------------------------------------------
    # Slide 5: Strategic Recommendations (White Background)
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_background(slide5, pure_white)
    add_header(slide5, "Action Recommendations")
    
    rec_box = slide5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(3.9))
    tf_rec = rec_box.text_frame
    tf_rec.word_wrap = True
    
    for idx, rec in enumerate(insights.get("recommendations", [])[:3]):
        p_rec = tf_rec.paragraphs[0] if idx == 0 else tf_rec.add_paragraph()
        p_rec.text = f"Priority {rec.get('priority')} - {rec.get('action')}"
        p_rec.font.name = "Segoe UI"
        p_rec.font.size = Pt(14)
        p_rec.font.bold = True
        p_rec.font.color.rgb = bright_blue
        p_rec.space_before = Pt(12)
        p_rec.space_after = Pt(3)
        
        p_rec_det = tf_rec.add_paragraph()
        p_rec_det.text = f"  • Rationale: {rec.get('rationale')}\n  • Expected Impact: {rec.get('expected_impact')}"
        p_rec_det.font.name = "Segoe UI"
        p_rec_det.font.size = Pt(10.5)
        p_rec_det.font.color.rgb = slate_navy
        
    # Speaker Notes
    slide5.notes_slide.notes_text_frame.text = (
        "These are our three primary strategic recommendations. "
        "Focus should be placed on high-priority items first. "
        "The expected impact is targeted towards reducing turnover or raising transactional AOV."
    )
    
    # ----------------------------------------------------
    # Slide 6: Meeting Talking Points (White Background)
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_background(slide6, pure_white)
    add_header(slide6, "Executive Briefing Script")
    
    tp_box = slide6.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9.0), Inches(3.9))
    tf_tp = tp_box.text_frame
    tf_tp.word_wrap = True
    
    for idx, tp in enumerate(insights.get("talking_points", [])[:5]):
        p_tp = tf_tp.paragraphs[0] if idx == 0 else tf_tp.add_paragraph()
        p_tp.text = f"{idx}. {tp}"
        p_tp.font.name = "Segoe UI"
        p_tp.font.size = Pt(13)
        p_tp.font.color.rgb = slate_navy
        p_tp.space_after = Pt(12)
        p_tp.line_spacing = 1.2
        
    # Speaker Notes
    slide6.notes_slide.notes_text_frame.text = (
        "This slide holds the core briefing script for the meeting. "
        "I will read these talking points in sequence to align the board."
    )
    
    # Save Presentation to Memory
    ppt_file = io.BytesIO()
    prs.save(ppt_file)
    ppt_file.seek(0)
    return ppt_file
